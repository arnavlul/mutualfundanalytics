import os
from reportlab.lib import colors
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, HRFlowable, Table, TableStyle
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

BASE_DIR = os.path.abspath(os.path.join(os.path.dirname(__file__), '..'))
reports_dir = os.path.join(BASE_DIR, 'reports')
os.makedirs(reports_dir, exist_ok=True)

def generate_pdf_report():
    pdf_path = os.path.join(reports_dir, 'Final_Report.pdf')
    doc = SimpleDocTemplate(pdf_path, pagesize=letter,
                            rightMargin=54, leftMargin=54, topMargin=54, bottomMargin=54)
    styles = getSampleStyleSheet()
    
    # Custom Styles
    title_style = ParagraphStyle(
        'DocTitle',
        parent=styles['Heading1'],
        fontName='Helvetica-Bold',
        fontSize=22,
        leading=26,
        textColor=colors.HexColor('#1B365D'),
        spaceAfter=12
    )
    subtitle_style = ParagraphStyle(
        'DocSubTitle',
        parent=styles['Normal'],
        fontName='Helvetica',
        fontSize=13,
        leading=16,
        textColor=colors.HexColor('#4A5568'),
        spaceAfter=20
    )
    h1_style = ParagraphStyle(
        'SectionHeading',
        parent=styles['Heading2'],
        fontName='Helvetica-Bold',
        fontSize=15,
        leading=18,
        textColor=colors.HexColor('#1B365D'),
        spaceBefore=14,
        spaceAfter=8
    )
    body_style = ParagraphStyle(
        'BodyTextCustom',
        parent=styles['Normal'],
        fontName='Helvetica',
        fontSize=10.5,
        leading=15,
        textColor=colors.HexColor('#2D3748'),
        spaceAfter=8
    )
    bullet_style = ParagraphStyle(
        'BulletCustom',
        parent=body_style,
        leftIndent=15,
        firstLineIndent=-10,
        spaceAfter=6
    )
    
    story = []
    
    # Title
    story.append(Paragraph("Bluestock Fintech: Mutual Fund Analytics Platform", title_style))
    story.append(Paragraph("<b>Capstone Project Final Report</b> — Data Engineering, ETL & Interactive Analytics", subtitle_style))
    story.append(HRFlowable(width="100%", thickness=2, color=colors.HexColor('#1B365D'), spaceBefore=0, spaceAfter=15))
    
    # 1. Executive Summary
    story.append(Paragraph("1. Executive Summary", h1_style))
    story.append(Paragraph(
        "The Indian mutual fund industry is experiencing rapid expansion, currently managing over ₹81 lakh crore in AUM across more than 1,900 schemes. However, retail investors and financial analysts frequently face data fragmentation across asset management companies (AMCs), making objective risk-adjusted comparisons challenging. This capstone project delivers a comprehensive, end-to-end analytics platform. By ingesting raw data from AMFI, mfapi.in, and market indices, our system cleans, transforms, and loads normalized data into an optimized SQLite database (`bluestock_mf.db`), culminating in interactive Power BI and Streamlit dashboards for real-time portfolio insights.",
        body_style
    ))
    
    # 2. Problem Statement & Objectives
    story.append(Paragraph("2. Problem Statement & Objectives", h1_style))
    story.append(Paragraph("<b>• Data Silos & Fragmentation:</b> Historical NAVs, AMC-level AUMs, SIP flows, and scheme details existed in isolated CSV and text formats. We unified them into a robust 5-table Star Schema.", bullet_style))
    story.append(Paragraph("<b>• Lack of Standardized Risk Metrics:</b> Investors lacked standardized tools to compute true risk-adjusted returns. We implemented automated calculation of 1/3/5-year CAGR, Sharpe Ratio, Sortino Ratio, Alpha, Beta, and Maximum Drawdown.", bullet_style))
    story.append(Paragraph("<b>• Benchmark Tracking Deficiencies:</b> We mapped scheme NAVs against NIFTY 50 and NIFTY 100 daily closing values to accurately compute tracking errors and market correlation.", bullet_style))
    story.append(Paragraph("<b>• Static Reporting Modernization:</b> Replaced legacy static reports with dynamic, self-serve interactive dashboards featuring multi-dimensional slicers and drill-downs.", bullet_style))
    
    # 3. Data Sources & ETL Architecture
    story.append(Paragraph("3. Data Sources & ETL Architecture", h1_style))
    story.append(Paragraph("Our automated data pipeline adheres to strict Data Engineering best practices:", body_style))
    story.append(Paragraph("<b>1. Extract (Ingestion):</b> Python scripts (`data_ingestion.py`, `live_nav_fetch.py`) extract raw data from AMFI daily feeds, mfapi.in historical REST APIs, and NSE/BSE index datasets into `data/raw/`.", bullet_style))
    story.append(Paragraph("<b>2. Transform (Cleaning):</b> Handled missing weekend and holiday NAV values using forward-filling (`ffill()`) across complete date reindexing. Standardized schema units to strictly distinguish between scheme-level AUM (in crores) and industry-wide AUM (in lakh crores).", bullet_style))
    story.append(Paragraph("<b>3. Load (Star Schema):</b> Loaded cleaned datasets into a normalized SQLite database (`data/db/bluestock_mf.db`) comprising `dim_fund`, `dim_date`, `fact_nav`, `fact_performance`, and `fact_transactions`.", bullet_style))
    story.append(Paragraph("<b>4. Analyze (Metrics):</b> Calculated advanced metrics including Historical VaR (95%), Conditional VaR (CVaR), and Sector Herfindahl-Hirschman Index (HHI) for concentration risk.", bullet_style))
    
    # 4. Key EDA & Performance Findings
    story.append(Paragraph("4. Key EDA & Performance Findings", h1_style))
    story.append(Paragraph("<b>• Record SIP Inflows:</b> Monthly SIP contributions hit historical highs, crossing ₹31,000+ crore, driven heavily by retail participation in mid and small-cap categories.", bullet_style))
    story.append(Paragraph("<b>• Risk-Return Tradeoff (Scorecard):</b> Our objective 0-100 Fund Scorecard revealed that while small-cap schemes delivered high raw returns, large-cap bluechip schemes offered superior Sortino ratios and significantly lower Maximum Drawdown during market volatility.", bullet_style))
    story.append(Paragraph("<b>• Sector Concentration (HHI):</b> Several thematic equity funds exhibited HHI scores exceeding 2,000, indicating severe over-concentration in Financial Services and IT sectors.", bullet_style))
    story.append(Paragraph("<b>• Cohort Dynamics:</b> The 2024 investor cohort represents the largest total AUM contribution, whereas the 2025 cohort displays higher transaction frequency with lower average ticket sizes.", bullet_style))
    
    # 5. Dashboard Implementation
    story.append(Paragraph("5. Dashboard Implementation & Visual Analytics", h1_style))
    story.append(Paragraph(
        "The visual presentation layer spans two robust platforms: an interactive Power BI dashboard (`dashboard/bluestock_mf.pbix`) and a Python-native Streamlit web application (`dashboard/app.py`). Every dashboard page features at least two functional slicers (Category, Risk Grade, and AMC selection), enabling stakeholders to drill into fund performance, analyze rolling 90-day Sharpe trends, and assess SIP continuity risk scores.",
        body_style
    ))
    
    # 6. Limitations & Actionable Recommendations
    story.append(Paragraph("6. Project Limitations & Actionable Recommendations", h1_style))
    story.append(Paragraph("<b>• Limitations:</b> Historical NAV feeds do not reflect intra-day price volatility or real-time expense ratio revisions. Additionally, investor transaction logs are synthetically modeled from aggregate geographic distributions (`T30` vs `B30` cities).", bullet_style))
    story.append(Paragraph("<b>• Recommendation 1 (Pipeline Automation):</b> Schedule the master `etl_pipeline.py` script as a daily automated cron job or via Apache Airflow to continuously refresh NAV histories.", bullet_style))
    story.append(Paragraph("<b>• Recommendation 2 (Portfolio Optimization):</b> Upgrade the standalone fund recommender to incorporate Markowitz Efficient Frontier modeling, allowing investors to build multi-fund portfolios along optimal risk-return boundaries.", bullet_style))
    
    doc.build(story)
    print("Generated Final_Report.pdf successfully!")

def generate_pptx_presentation():
    pptx_path = os.path.join(reports_dir, 'Presentation.pptx')
    prs = Presentation()
    
    # Use 16:9 widescreen slides
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    def add_slide_content(title_text, bullets):
        slide = prs.slides.add_slide(prs.slide_layouts[1]) # Title and Content layout
        title = slide.shapes.title
        title.text = title_text
        title.text_frame.paragraphs[0].font.name = 'Arial'
        title.text_frame.paragraphs[0].font.size = Pt(28)
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(27, 54, 93)
        
        body = slide.placeholders[1]
        tf = body.text_frame
        tf.clear()
        
        for i, bullet in enumerate(bullets):
            p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
            p.text = bullet
            p.font.name = 'Arial'
            p.font.size = Pt(18)
            p.font.color.rgb = RGBColor(45, 55, 72)
            p.space_after = Pt(14)
            p.level = 0
        return slide

    # Slide 1: Title
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Mutual Fund Analytics Platform"
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.size = Pt(40)
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(27, 54, 93)
    subtitle.text = "End-to-End Data Engineering, ETL Pipeline & Interactive Dashboard\nBluestock Fintech Capstone Submission"
    subtitle.text_frame.paragraphs[0].font.size = Pt(22)

    # Slide 2: Problem Statement
    add_slide_content("Problem Statement & Challenges", [
        "Data Fragmentation: NAV histories, AMC AUMs, and SIP data existed in isolated silos.",
        "Lack of Standardization: No automated way to compare risk-adjusted metrics across 1,900+ schemes.",
        "Benchmark Disconnect: Difficult to measure true scheme Alpha vs NIFTY 50 and NIFTY 100 indices.",
        "Legacy Reporting: Static monthly PDFs lacked interactive filtering and drill-down capabilities."
    ])

    # Slide 3: Project Objectives
    add_slide_content("Project Objectives & Deliverables", [
        "Build an End-to-End ETL Pipeline: Extract from public APIs/CSVs, clean, and normalize data.",
        "Normalize Database Architecture: Implement an optimized 5-table Star Schema in SQLite.",
        "Automate Risk & Performance Metrics: Compute CAGR, Sharpe, Sortino, Alpha, Beta, and Max Drawdown.",
        "Deliver Interactive Visualizations: Build multi-page Power BI and Streamlit dashboards."
    ])

    # Slide 4: Data Sources & Volume
    add_slide_content("Data Sources & Ingestion Scope", [
        "AMFI India & mfapi.in: Daily NAV updates, scheme master lists, and historical price feeds.",
        "NSE / BSE Benchmarks: Historical daily closing values for NIFTY 50 and NIFTY 100 indices.",
        "Volume Processed: Over 46,000+ NAV history records and 87,000+ transaction logs.",
        "Automated Ingestion: Python scripts fetch, validate, and store raw datasets in 'data/raw/'."
    ])

    # Slide 5: System & ETL Architecture
    add_slide_content("ETL & Database Architecture (Star Schema)", [
        "Extract: Python ingestion modules (`data_ingestion.py`, `live_nav_fetch.py`).",
        "Transform: Forward-filling (`ffill()`) missing holiday NAVs; explicit unit separation (crore vs lakh crore).",
        "Load: SQLite database (`data/db/bluestock_mf.db`) with relational constraints.",
        "Schema Design: 2 Dimension Tables (`dim_fund`, `dim_date`) and 3 Fact Tables (`fact_nav`, `fact_performance`, `fact_transactions`)."
    ])

    # Slide 6: EDA Highlights (Market & SIP Trends)
    add_slide_content("EDA Highlights: Market & SIP Growth", [
        "Record SIP Contributions: Monthly SIP inflows crossed ₹31,000 crore, reflecting robust retail adoption.",
        "AMC Concentration: Top 5 AMCs (e.g., SBI, ICICI Prudential, HDFC) control nearly 50% of total industry AUM.",
        "Category Preferences: Strong shift in retail folio counts toward Mid-Cap and Small-Cap equity schemes.",
        "Demographic Split: Significant transaction volume variance between T30 (Top 30 cities) and B30 regions."
    ])

    # Slide 7: Performance & Risk Analytics
    add_slide_content("Performance & Risk Calculation Methodology", [
        "Standardized Annualization: Strictly used 252 trading days (`252 / n_trading_days`) for all CAGR models.",
        "Sharpe & Sortino Ratios: Evaluated total risk-adjusted returns versus downside-only volatility.",
        "Alpha & Beta Regression: Linear regression against NIFTY 100 daily returns to isolate manager skill.",
        "Drawdown Analysis: Tracked running maximums to identify historical peak-to-trough drops."
    ])

    # Slide 8: The 0-100 Fund Scorecard
    add_slide_content("The Objective Fund Scorecard Model", [
        "Composite Ranking System: Graded funds out of 100 using a balanced multi-metric weighted formula.",
        "Weighting Breakdown: 3yr CAGR (30%), Sharpe Ratio (25%), Alpha (20%), Expense Ratio (15%), Max Drawdown (10%).",
        "Key Takeaway: Bluechip Large-Cap schemes consistently outscore high-volatility Small-Caps on risk-adjusted terms.",
        "Exported Output: Automatically generated and saved to `data/processed/fund_scorecard.csv`."
    ])

    # Slide 9: Advanced Analytics (VaR & HHI)
    add_slide_content("Advanced Analytics: VaR, CVaR & Sector HHI", [
        "Historical VaR & CVaR (95%): Quantified exact daily loss thresholds and expected shortfall in worst-case scenarios.",
        "Rolling 90-Day Sharpe: Identified periods of performance instability during broader market corrections.",
        "Sector Concentration (HHI): Calculated Herfindahl-Hirschman Index across equity holdings.",
        "Concentration Alerts: Schemes with HHI > 2,000 highlighted for extreme overweight positions in Banking/IT."
    ])

    # Slide 10: Interactive Dashboard Features
    add_slide_content("Interactive Dashboard Capabilities", [
        "Multi-Platform Delivery: Built in both Power BI (`bluestock_mf.pbix`) and Python Streamlit (`app.py`).",
        "Four Core Pages: Industry Overview, Fund Performance & Scorecard, Advanced Risk, and Investor Demographics.",
        "Slicer Interactivity: Every page equipped with at least 2 functional slicers (Category, Risk Grade, AMC).",
        "Self-Serve Exploration: Users can filter out high-risk funds with a single click."
    ])

    # Slide 11: Actionable Recommendations
    add_slide_content("Actionable Recommendations & Next Steps", [
        "Automate Daily Cron Jobs: Schedule `etl_pipeline.py` to auto-refresh NAVs from mfapi.in every weekday at 8 PM.",
        "Modern Portfolio Theory: Expand our standalone fund recommender to compute Markowitz Efficient Frontiers.",
        "Real-Time Alerting: Integrate automated email alerts when a scheme's rolling Sharpe ratio drops below benchmark.",
        "Cloud Deployment: Migrate local SQLite star schema to AWS RDS (PostgreSQL) for enterprise multi-user scaling."
    ])

    # Slide 12: Conclusion & Q&A
    add_slide_content("Summary & Conclusion", [
        "Unified Platform: Successfully bridged data silos into a clean, reproducible engineering workflow.",
        "High-Fidelity Analytics: Delivered objective scoring, advanced VaR modeling, and cohort insights.",
        "Complete Deliverables: Code, schema, notebooks, and reports structured precisely to capstone standards.",
        "Thank You — Open for Questions & Feedback!"
    ])

    prs.save(pptx_path)
    print("Generated Presentation.pptx successfully!")

if __name__ == '__main__':
    generate_pdf_report()
    generate_pptx_presentation()
